#!/usr/bin/env python3
"""
generate_report_pptx.py - render the branded competitor-analysis as a PPTX deck.

Same content + brand styling as generate_report.py (the PDF), with two differences:
  - Output is a .pptx (python-pptx). Slides are PORTRAIT 8.5x11in, i.e. the exact aspect
    ratio of the PDF's LETTER page (the deliverable is a deck, not a wide 16:9 presentation).
  - Sources are attributed *per paragraph and table* (comprehensive), not only listed at the end:
      * market-wide text (exec summary, trends, what's working, opportunities) -> the global source list
      * each comparison table (landscape, marketing) -> the union of the competitors' sources
      * each competitor profile -> that competitor's own sources
    A full numbered bibliography still closes the deck.

Self-contained: brand kit, the strengths/watch-outs chart (matplotlib), and run-data loading all live here.

Reads:  clients/<slug>/data/runs/<date>/{snapshot.json, changes.json}, clients/<slug>/brand_kit.json
Writes: clients/<slug>/reports/<slug>_<date>_competitor_report.pptx

CLI:
    uv run python tools/generate_report_pptx.py [--client ivp] [--date 2026-06-15] [--out path.pptx]
"""
from __future__ import annotations

import argparse
import datetime
import json
import math
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import common  # noqa: E402

import matplotlib  # noqa: E402
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402


# ---------------------------------------------------------------- geometry (8.5x11 = LETTER aspect)
PAGE_W = Inches(8.5)
PAGE_H = Inches(11)
ML = Inches(0.85)          # left margin (matches the PDF)
MR = Inches(0.85)
CONTENT_W = Inches(6.8)    # 8.5 - 0.85 - 0.85
CW_IN = 6.8                # content width in inches, for text-height estimation
TOP_Y = Inches(1.0)        # first content baseline (below the running header)
BOTTOM_Y = Inches(10.35)   # content must stop above the footer

DARK = "#1a1f29"           # body text colour (overridden per brand in build_pptx)
MUTED = "#5b6470"
LABEL = "#6b7280"


# ---------------------------------------------------------------- run data + chart
def latest_run(runs_dir):
    runs = [d for d in os.listdir(runs_dir)
            if os.path.isfile(os.path.join(runs_dir, d, "snapshot.json"))] if os.path.isdir(runs_dir) else []
    return sorted(runs)[-1] if runs else None


def load(path, default=None):
    if path and os.path.exists(path):
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    return default


def build_chart(competitors, C, path):
    """Horizontal bar chart: strengths vs. watch-outs identified per competitor."""
    comps = [c for c in competitors if (c.get("strengths") or c.get("weaknesses"))]
    if not comps:
        return None
    names = [c["name"] for c in comps]
    strengths = [len(c.get("strengths") or []) for c in comps]
    weaknesses = [-len(c.get("weaknesses") or []) for c in comps]
    y = range(len(names))
    fig, ax = plt.subplots(figsize=(7.0, 0.5 * len(names) + 1.2), dpi=150)
    ax.barh(y, strengths, color=C["accent"], label="Strengths", height=0.6)
    ax.barh(y, weaknesses, color=C["secondary"], label="Watch-outs", height=0.6)
    ax.set_yticks(list(y))
    ax.set_yticklabels(names, fontsize=9)
    ax.axvline(0, color="#9aa0a6", linewidth=0.8)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(length=0)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.12), ncol=2, fontsize=8, frameon=False)
    ax.set_title("Strengths vs. watch-outs identified per competitor", fontsize=11,
                 color=C["primary"], fontweight="bold", loc="left")
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


# ---------------------------------------------------------------- small helpers
def RGB(hexstr: str) -> RGBColor:
    return RGBColor.from_string((hexstr or "#000000").lstrip("#"))


def brand_font(brand: dict) -> str:
    return (brand.get("fonts") or {}).get("family_name") or "Arial"


def dedupe(seq):
    seen, out = set(), []
    for s in seq or []:
        s = (s or "").strip() if isinstance(s, str) else s
        if s and s not in seen:
            seen.add(s)
            out.append(s)
    return out


def est_lines(text: str, size: float, width_in: float) -> int:
    """Rough number of wrapped lines for proportional text at a given pt size."""
    char_w_in = size * 0.52 / 72.0          # avg glyph ~0.52em, generous to avoid overlap
    cpl = max(1, int(width_in / char_w_in))
    return max(1, math.ceil(len(str(text)) / cpl))


def est_height(text: str, size: float, width_in: float, leading: float = 1.32) -> Emu:
    return Inches(est_lines(text, size, width_in) * size * leading / 72.0)


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def _run(p, text, *, size, color, bold=False, italic=False, font=None):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = RGB(color)
    if font:
        f.name = font
    return r


def rect(slide, left, top, width, height, fill_hex):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGB(fill_hex)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


# ---------------------------------------------------------------- flow placement (top-down y cursor)
def place_para(slide, y, text=None, *, left=ML, width=CONTENT_W, width_in=CW_IN, size=10,
               color=None, bold=False, italic=False, font=None, align=PP_ALIGN.LEFT,
               gap=6.0, runs=None):
    color = color or DARK
    measure = text if runs is None else "".join(t for t, _ in runs)
    h = est_height(measure, size, width_in)
    tf = textbox(slide, left, y, width, h)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    if runs is None:
        _run(p, text, size=size, color=color, bold=bold, italic=italic, font=font)
    else:
        for t, o in runs:
            _run(p, t, size=o.get("size", size), color=o.get("color", color),
                 bold=o.get("bold", False), italic=o.get("italic", False), font=o.get("font", font))
    return y + h + Pt(gap)


def place_bullets(slide, y, items, *, left=ML, width=CONTENT_W, width_in=CW_IN, size=9.5,
                  color=None, font=None, gap=6.0, limit=None, glyph="•  "):
    color = color or DARK
    items = [str(i).strip() for i in (items or []) if str(i).strip()]
    if limit:
        items = items[:limit]
    if not items:
        return place_para(slide, y, "None identified.", left=left, width=width, width_in=width_in,
                          size=size - 1, color=MUTED, italic=True, font=font, gap=gap)
    total_lines = sum(est_lines(glyph + it, size, width_in - 0.1) for it in items)
    h = Inches(total_lines * size * 1.32 / 72.0) + Pt(3 * len(items))
    tf = textbox(slide, left, y, width, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        _run(p, glyph + it, size=size, color=color, font=font)
    return y + h + Pt(gap)


def label_para(slide, y, label, value, *, left=ML, width=CONTENT_W, width_in=CW_IN,
               value_size=8.5, font=None, gap=3.0):
    return place_para(slide, y, left=left, width=width, width_in=width_in, size=value_size, font=font, gap=gap,
                      runs=[(f"{label.upper()}  ", {"bold": True, "color": LABEL, "size": value_size - 1}),
                            (str(value), {"size": value_size, "color": DARK})])


def sources_line(slide, y, urls, *, font=None, label="Sources", limit=8):
    urls = dedupe(urls)
    if not urls:
        return y
    shown = urls[:limit]
    extra = len(urls) - len(shown)
    text = f"{label}: " + "   ·   ".join(shown) + (f"   (+{extra} more — see Sources & Methodology)" if extra else "")
    return place_para(slide, y, text, size=7, color=LABEL, font=font, gap=4)


# ---------------------------------------------------------------- tables
def _plain_table(tbl):
    """Strip the default theme table style so our explicit fills/fonts render cleanly."""
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is None:
        return
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    for el in list(tblPr.findall(qn("a:tableStyleId"))):
        tblPr.remove(el)
    sid = tblPr.makeelement(qn("a:tableStyleId"), {})
    sid.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # No Style, Table Grid
    tblPr.append(sid)


def _cell(cell, text, *, size, color, fill, bold=False, font=None, align=PP_ALIGN.LEFT):
    cell.margin_left = Pt(5)
    cell.margin_right = Pt(5)
    cell.margin_top = Pt(3)
    cell.margin_bottom = Pt(3)
    cell.vertical_anchor = MSO_ANCHOR.TOP
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGB(fill)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, str(text), size=size, color=color, bold=bold, font=font)


def build_table(slide, y, header, rows, col_w_in, ctx, *, body_size=9, header_size=8.5, left=ML):
    C, font = ctx["C"], ctx["font"]
    n_rows, n_cols = len(rows) + 1, len(header)

    def row_h_in(cells, size):
        mx = max((est_lines(c, size, w - 0.14) for c, w in zip(cells, col_w_in)), default=1)
        return mx * size * 1.32 / 72.0 + 0.11

    hh = row_h_in(header, header_size)
    body_hs = [row_h_in(r, body_size) for r in rows]
    total_h_in = hh + sum(body_hs)

    shape = slide.shapes.add_table(n_rows, n_cols, left, y, Inches(sum(col_w_in)), Inches(total_h_in))
    tbl = shape.table
    _plain_table(tbl)
    for i, w in enumerate(col_w_in):
        tbl.columns[i].width = Inches(w)
    tbl.rows[0].height = Inches(hh)
    for i, h in enumerate(body_hs):
        tbl.rows[i + 1].height = Inches(h)
    for ci, htext in enumerate(header):
        _cell(tbl.cell(0, ci), htext, size=header_size, color="#FFFFFF", fill=C["primary"], bold=True, font=font)
    for ri, row in enumerate(rows):
        fill = "#FFFFFF" if ri % 2 == 0 else C["light"]
        for ci, val in enumerate(row):
            _cell(tbl.cell(ri + 1, ci), val, size=body_size, color=DARK, fill=fill, bold=(ci == 0), font=font)
    return y + Inches(total_h_in) + Pt(8)


# ---------------------------------------------------------------- page furniture
def heading(slide, title, ctx, y):
    C, font = ctx["C"], ctx["font"]
    rect(slide, ML, y, Inches(1.3), Pt(3), C["accent"])
    tf = textbox(slide, ML, y + Pt(7), CONTENT_W, Inches(0.5))
    _run(tf.paragraphs[0], title, size=17, color=C["primary"], bold=True, font=font)
    return y + Pt(7) + Inches(0.42)


def add_chrome(slide, ctx):
    C, font = ctx["C"], ctx["font"]
    x = ML
    icon = ctx["logo_icon"]
    if icon and os.path.exists(icon):
        try:
            slide.shapes.add_picture(icon, ML, Inches(0.45), height=Inches(0.26))
            x = ML + Inches(0.36)
        except Exception:
            pass
    tf = textbox(slide, x, Inches(0.46), Inches(4.5), Inches(0.3))
    _run(tf.paragraphs[0], ctx["title"], size=9, color=C["primary"], bold=True, font=font)
    rect(slide, ML, Inches(0.8), PAGE_W - ML - MR, Pt(1.2), C["accent"])
    rect(slide, ML, PAGE_H - Inches(0.62), PAGE_W - ML - MR, Pt(0.6), "#e5e7eb")
    tf = textbox(slide, ML, PAGE_H - Inches(0.55), Inches(4.5), Inches(0.3))
    _run(tf.paragraphs[0], f"Competitive Landscape Report · {ctx['date']}", size=8, color=LABEL, font=font)
    tf = textbox(slide, PAGE_W - MR - Inches(1.6), PAGE_H - Inches(0.55), Inches(1.6), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _run(p, f"Page {ctx['page']}", size=8, color=LABEL, font=font)


def content_slide(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    ctx["page"] += 1
    add_chrome(slide, ctx)
    return slide


def build_cover(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    C, font = ctx["C"], ctx["font"]
    rect(slide, 0, 0, PAGE_W, PAGE_H, C["primary"])
    rect(slide, 0, Inches(0.55), PAGE_W, Inches(0.12), C["accent"])
    logo = ctx["logo_white"]
    if logo and os.path.exists(logo):
        try:
            disp_w = Inches(2.8)
            slide.shapes.add_picture(logo, int((PAGE_W - disp_w) / 2), Inches(2.1), width=disp_w)
        except Exception:
            pass
    tf = textbox(slide, Inches(0.5), Inches(4.6), PAGE_W - Inches(1.0), Inches(1.0))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "Competitive Landscape Report", size=30, color="#FFFFFF", bold=True, font=font)
    tf = textbox(slide, Inches(0.5), Inches(5.55), PAGE_W - Inches(1.0), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, ctx["subtitle"], size=14, color=C["accent"], font=font)
    tf = textbox(slide, Inches(0.5), Inches(6.05), PAGE_W - Inches(1.0), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, ctx["date"], size=11, color="#c7ccd6", font=font)
    tf = textbox(slide, Inches(0.5), PAGE_H - Inches(0.85), PAGE_W - Inches(1.0), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, f"Confidential — prepared for {ctx['title']}", size=8.5, color="#8b93a3", font=font)


# ---------------------------------------------------------------- competitor slide
def build_competitor_slide(prs, ctx, c):
    C, font = ctx["C"], ctx["font"]
    slide = content_slide(prs, ctx)
    y = heading(slide, c.get("name", "Competitor"), ctx, TOP_Y)
    if c.get("one_liner"):
        y = place_para(slide, y, c["one_liner"], size=10.5, italic=True, color=C["primary"], font=font, gap=7)

    col_gap = 0.3
    colw = (CW_IN - col_gap) / 2.0           # ~3.25in per column
    right_left = ML + Inches(colw + col_gap)

    # left column: facts
    yl = y
    for label, val in [("Positioning", c.get("positioning")),
                       ("Target customers", c.get("target_customers")),
                       ("Products", ", ".join(c.get("products") or [])),
                       ("Key features", ", ".join((c.get("key_features") or [])[:5])),
                       ("Pricing", c.get("pricing")),
                       ("Ratings", c.get("ratings")),
                       ("Sentiment", c.get("customer_sentiment"))]:
        if str(val or "").strip():
            yl = label_para(slide, yl, label, val, left=ML, width=Inches(colw), width_in=colw, value_size=8.5, font=font)

    # right column: strengths / watch-outs
    yr = y
    yr = place_para(slide, yr, "STRENGTHS", left=right_left, width=Inches(colw), width_in=colw,
                    size=8, bold=True, color=LABEL, font=font, gap=2)
    yr = place_bullets(slide, yr, c.get("strengths"), left=right_left, width=Inches(colw), width_in=colw,
                       size=8.5, font=font, limit=5, gap=5)
    yr = place_para(slide, yr, "WATCH-OUTS", left=right_left, width=Inches(colw), width_in=colw,
                    size=8, bold=True, color=LABEL, font=font, gap=2)
    yr = place_bullets(slide, yr, c.get("weaknesses"), left=right_left, width=Inches(colw), width_in=colw,
                       size=8.5, font=font, limit=5, gap=5)

    y = max(yl, yr) + Pt(4)
    if c.get("recent_moves"):
        y = place_para(slide, y, "RECENT MOVES", size=8, bold=True, color=LABEL, font=font, gap=2)
        y = place_bullets(slide, y, c.get("recent_moves"), size=8.5, font=font, limit=4, gap=5)
    mk = c.get("marketing_content") or c.get("content_marketing")
    if mk and y < BOTTOM_Y - Inches(1.0):
        y = label_para(slide, y, "Marketing & channels", mk, value_size=8.5, font=font, gap=5)
    sources_line(slide, y, c.get("sources") or [], font=font)


# ---------------------------------------------------------------- main build
def build_pptx(snapshot, changes, brand, out_path):
    global DARK
    C = brand["colors"]
    DARK = C.get("dark", "#1a1f29")
    font = brand_font(brand)
    biz = snapshot.get("business", {})
    date_str = snapshot.get("run_date", datetime.date.today().isoformat())
    title = biz.get("name", "Your Company")
    comps = snapshot.get("competitors", [])
    ms = snapshot.get("market_summary", {})
    logo = brand.get("logo") or {}
    ctx = {"C": C, "font": font, "title": title,
           "subtitle": biz.get("tagline") or biz.get("category") or "Market & Competitor Intelligence",
           "date": date_str, "logo_white": logo.get("primary"), "logo_icon": logo.get("icon"), "page": 0}

    global_sources = dedupe(snapshot.get("sources"))
    comp_union = dedupe([s for c in comps for s in (c.get("sources") or [])])

    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    build_cover(prs, ctx)

    # Executive summary
    s = content_slide(prs, ctx)
    y = heading(s, "Executive Summary", ctx, TOP_Y)
    if ms.get("overview"):
        y = place_para(s, y, ms["overview"], size=10, font=font, gap=6)
    y = place_para(s, y, font=font, gap=6,
                   runs=[(f"{len(comps)} competitors tracked: ", {"bold": True, "color": C["primary"]}),
                         (", ".join(c["name"] for c in comps) + ".", {})])
    if changes and changes.get("summary") and not changes.get("baseline"):
        y = place_para(s, y, "Headline changes since last run", size=8.5, bold=True, color=LABEL, font=font, gap=2)
        y = place_bullets(s, y, changes["summary"], size=9.5, font=font, limit=5)
    if ms.get("key_trends"):
        y = place_para(s, y, "Key market trends", size=8.5, bold=True, color=LABEL, font=font, gap=2)
        y = place_bullets(s, y, ms["key_trends"], size=9.5, font=font, limit=6)
    sources_line(s, y, global_sources, font=font)

    # Your business
    s = content_slide(prs, ctx)
    y = heading(s, f"Your Business — {title}", ctx, TOP_Y)
    diff = biz.get("differentiators")
    rows = [("What you do", biz.get("description") or biz.get("value_prop")),
            ("Category", biz.get("category")),
            ("Target customers", biz.get("target_customers")),
            ("Value proposition", biz.get("value_prop")),
            ("Differentiators", ", ".join(diff) if isinstance(diff, list) else diff),
            ("Website", biz.get("url"))]
    rows = [[k.upper(), v] for k, v in rows if str(v or "").strip()]
    y = build_table(s, y, ["FIELD", "DETAIL"], rows, [1.6, 5.2], ctx, body_size=9)
    sources_line(s, y, [biz.get("url")], font=font)

    # Market landscape
    s = content_slide(prs, ctx)
    y = heading(s, "Market Landscape", ctx, TOP_Y)
    rows = [[c["name"], c.get("one_liner") or c.get("positioning") or ""] for c in comps]
    y = build_table(s, y, ["COMPETITOR", "POSITIONING"], rows, [1.7, 5.1], ctx, body_size=9)
    sources_line(s, y, comp_union, font=font)

    # Strengths chart
    chart = build_chart(comps, C, os.path.join(".tmp", "chart_strengths.png"))
    if chart and os.path.exists(chart):
        s = content_slide(prs, ctx)
        y = heading(s, "Competitive Strength Map", ctx, TOP_Y)
        pic = s.shapes.add_picture(chart, ML, y, width=Inches(6.6))
        sources_line(s, y + pic.height + Pt(10), comp_union, font=font)

    # Marketing & channel presence
    if any(c.get("social_presence") or c.get("content_focus") for c in comps):
        s = content_slide(prs, ctx)
        y = heading(s, "Marketing & Channel Presence", ctx, TOP_Y)
        rows = []
        for c in comps:
            sp = c.get("social_presence") or {}
            rows.append([c["name"], sp.get("linkedin") or "—", c.get("youtube_activity") or "—",
                         c.get("content_focus") or ""])
        y = build_table(s, y, ["COMPETITOR", "LINKEDIN", "YOUTUBE", "CONTENT FOCUS"], rows,
                        [1.5, 1.5, 1.6, 2.2], ctx, body_size=8.5)
        note = ms.get("channel_gap_note") or ms.get("ivp_channel_note")
        if note:
            y = place_para(s, y, font=font, gap=6,
                           runs=[("Your channel gap: ", {"bold": True, "color": C["primary"]}), (note, {})])
        sources_line(s, y, comp_union, font=font)

    # Per-competitor profiles (one slide each)
    for c in comps:
        build_competitor_slide(prs, ctx, c)

    # What's working for competitors
    if ms.get("whats_working_for_competitors"):
        s = content_slide(prs, ctx)
        y = heading(s, "What's Working for Competitors", ctx, TOP_Y)
        y = place_bullets(s, y, ms["whats_working_for_competitors"], size=10, font=font)
        sources_line(s, y, comp_union, font=font)

    # Where you can improve
    opps = ms.get("your_opportunities") or []
    if opps:
        s = content_slide(prs, ctx)
        y = heading(s, "Where You Can Improve", ctx, TOP_Y)
        for o in opps:
            if y > BOTTOM_Y - Inches(1.1):
                s = content_slide(prs, ctx)
                y = heading(s, "Where You Can Improve (cont.)", ctx, TOP_Y)
            if isinstance(o, dict):
                pr = (o.get("priority") or "").upper()
                r = [(o.get("title", ""), {"bold": True, "color": C["primary"], "size": 10.5})]
                if pr:
                    r.append((f"   [{pr}]", {"bold": True, "color": C["accent"], "size": 8.5}))
                y = place_para(s, y, runs=r, font=font, gap=2)
                if o.get("rationale"):
                    y = place_para(s, y, o["rationale"], size=9, color=MUTED, font=font, gap=7)
            else:
                y = place_para(s, y, "•  " + str(o), size=10, font=font, gap=5)
        sources_line(s, y, global_sources, font=font)

    # What changed since last run
    s = content_slide(prs, ctx)
    y = heading(s, "What Changed Since Last Run", ctx, TOP_Y)
    if changes and changes.get("baseline"):
        y = place_para(s, y, "This is the baseline run — the first snapshot of the market. "
                             "Future runs will highlight what changed here.", size=10, font=font)
    elif changes:
        y = place_para(s, y, f"Compared against the run on {changes.get('previous_date')}.",
                       size=9, color=MUTED, font=font, gap=4)
        y = place_bullets(s, y, changes.get("summary"), size=10, font=font)
    else:
        y = place_para(s, y, "No change data available for this run.", size=9, color=MUTED, font=font)

    # Sources & methodology (full numbered bibliography)
    s = content_slide(prs, ctx)
    y = heading(s, "Sources & Methodology", ctx, TOP_Y)
    if snapshot.get("methodology_note"):
        y = place_para(s, y, snapshot["methodology_note"], size=8.5, color=MUTED, font=font, gap=8)
    all_sources = dedupe((snapshot.get("sources") or []) + comp_union)
    for i, u in enumerate(all_sources, 1):
        if y > BOTTOM_Y:
            s = content_slide(prs, ctx)
            y = heading(s, "Sources & Methodology (cont.)", ctx, TOP_Y)
        y = place_para(s, y, f"{i}.  {u}", size=8, color=MUTED, font=font, gap=2)

    prs.save(out_path)
    return out_path, len(prs.slides._sldIdLst)


def main(argv=None):
    ap = argparse.ArgumentParser(description="Render the branded competitor-analysis as a PPTX deck.")
    ap.add_argument("--date", default=None, help="Run date to render (default: latest)")
    common.add_client_arg(ap)
    ap.add_argument("--out", default=None, help="Output .pptx path (default: clients/<slug>/reports/...)")
    args = ap.parse_args(argv)

    slug = common.resolve_slug(args.client)
    paths = common.client_paths(slug)
    date = args.date or latest_run(paths["runs"])
    if not date:
        print(f"No snapshot found to render for client '{slug}'.", file=sys.stderr)
        return 1
    run_dir = os.path.join(paths["runs"], date)
    snapshot = load(os.path.join(run_dir, "snapshot.json"))
    if not snapshot:
        print(f"No snapshot.json in {run_dir}", file=sys.stderr)
        return 1
    changes = load(os.path.join(run_dir, "changes.json"))
    brand = load(paths["brand_kit"])
    if not brand:
        print(f"No brand kit at {paths['brand_kit']}", file=sys.stderr)
        return 1

    os.makedirs(paths["reports"], exist_ok=True)
    os.makedirs(".tmp", exist_ok=True)
    out = args.out or os.path.join(paths["reports"], f"{slug}_{date}_competitor_report.pptx")
    out, n_slides = build_pptx(snapshot, changes, brand, out)
    size = os.path.getsize(out)
    print(f"Wrote {out}  ({size // 1024} KB, {len(comps) if (comps := snapshot.get('competitors', [])) else 0} "
          f"competitors, {n_slides} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
